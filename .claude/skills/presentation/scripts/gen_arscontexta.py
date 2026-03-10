#!/usr/bin/env python3
"""Generate arscontexta presentation with custom layouts per slide."""
import sys
sys.path.insert(0, "/Users/viktorsolomonik/Projects/svaib/.venv/lib/python3.14/site-packages")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Brand Constants ───────────────────────────────────────────
TEAL = RGBColor(0x00, 0x8B, 0x7F)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
PINK = RGBColor(0xFF, 0x4D, 0x8D)
GRAY = RGBColor(0x6B, 0x72, 0x80)
BRIGHT_TEAL = RGBColor(0x00, 0xB4, 0xA6)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL_BG = RGBColor(0xF0, 0xFD, 0xFB)
SUBTLE_TEAL = RGBColor(0xE0, 0xF7, 0xF5)
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)

FONT_H = "Montserrat"
FONT_B = "Roboto"

# ─── Helpers ───────────────────────────────────────────────────

def brand_markers(slide):
    """Add teal stripe + svaib footer to every slide."""
    # Teal stripe top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Footer: svaib.
    tf = slide.shapes.add_textbox(
        Inches(8.8), Inches(5.25), Inches(1.0), Inches(0.3)
    ).text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "svaib"
    run.font.name = FONT_H
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY
    run2 = p.add_run()
    run2.text = "."
    run2.font.name = FONT_H
    run2.font.size = Pt(8)
    run2.font.color.rgb = PINK


def add_title(slide, text, x=0.5, y=0.3, w=9.0, h=0.7, size=38, color=TEAL):
    """Add slide title."""
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_H
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return tf


def add_title_two_colors(slide, parts, x=0.5, y=0.3, w=9.0, h=0.7, size=38):
    """Add title with multiple colored parts. parts = [(text, color), ...]"""
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for text, color in parts:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_H
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
    return tf


def add_subtitle(slide, text, x=0.5, y=1.0, w=9.0, h=0.4, size=18, color=GRAY):
    """Add subtitle / description line."""
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_B
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return tf


def add_text(slide, text, x, y, w, h, font=None, size=16, color=DARK_BLUE,
             bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """Add a text box."""
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font or FONT_B
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def add_multiline(slide, lines, x, y, w, h, size=16, color=DARK_BLUE, spacing=8):
    """Add text with multiple paragraphs."""
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(line, tuple):
            text, lcolor, lbold = line[0], line[1], line[2] if len(line) > 2 else False
        else:
            text, lcolor, lbold = line, color, False
        run = p.add_run()
        run.text = text
        run.font.name = FONT_B
        run.font.size = Pt(size)
        run.font.color.rgb = lcolor
        run.font.bold = lbold
    return tf


def add_card(slide, x, y, w, h, fill=WHITE, border=CARD_BORDER, radius=0.1):
    """Add a rounded card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, x, y, w, h, fill=TEAL):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_circle(slide, x, y, size, fill=BRIGHT_TEAL):
    """Add a filled circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_footer_line(slide, text, y=4.85):
    """Add a footer summary line."""
    tf = slide.shapes.add_textbox(
        Inches(0.5), Inches(y), Inches(9.0), Inches(0.45)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_B
    run.font.size = Pt(15)
    run.font.color.rgb = GRAY
    return tf


# ─── Slide 1: Title ───────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    # Decorative teal block top-right
    add_rect(slide, 7.8, 0.6, 1.8, 1.6, SUBTLE_TEAL)

    # Title: "arscontexta"
    add_text(slide, "arscontexta", 1.0, 1.5, 8.0, 1.0,
             font=FONT_H, size=56, color=TEAL, bold=True)

    # Subtitle line 1
    add_text(slide, "как устроен context engineering", 1.0, 2.6, 8.0, 0.8,
             font=FONT_B, size=32, color=DARK_BLUE)

    # Decorative elements
    add_circle(slide, 8.3, 3.2, 0.7, BRIGHT_TEAL)
    add_circle(slide, 9.1, 3.8, 0.35, YELLOW)

    # Thin teal accent line under title
    add_rect(slide, 1.0, 2.45, 2.5, 0.04, TEAL)


# ─── Slide 2: Three Spaces ────────────────────────────────────

def slide_three_spaces(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title(slide, "Три пространства")
    add_subtitle(slide, "Архитектурный инвариант: не папки, а разные скорости жизни")

    # Three cards
    cards = [
        ("self/", "кто я", TEAL,
         ["Идентичность агента.", "identity · methodology · goals",
          "Загружается целиком"]),
        ("notes/", "что знаю", PINK,
         ["Ядро знаний.", "Атомарные файлы + навигаторы",
          "Загружается выборочно"]),
        ("ops/", "что делаю", BRIGHT_TEAL,
         ["Операционка.", "Задачи · логи · отчёты",
          "Загружается точечно"]),
    ]

    card_w = 2.7
    gap = 0.35
    start_x = 0.5 + (9.0 - 3 * card_w - 2 * gap) / 2

    for i, (name, label, accent, lines) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        y_card = 1.65

        # Card bg
        add_card(slide, x, y_card, card_w, 3.3, WHITE, CARD_BORDER)

        # Accent bar top of card
        add_rect(slide, x, y_card, card_w, 0.06, accent)

        # Folder name
        add_text(slide, name, x + 0.2, y_card + 0.25, card_w - 0.4, 0.5,
                 font=FONT_H, size=26, color=DARK_BLUE, bold=True)

        # Label in accent color
        add_text(slide, label, x + 0.2, y_card + 0.75, card_w - 0.4, 0.35,
                 font=FONT_B, size=18, color=accent)

        # Description lines
        for j, line in enumerate(lines):
            add_text(slide, line, x + 0.2, y_card + 1.25 + j * 0.45, card_w - 0.4, 0.4,
                     size=14, color=DARK_BLUE)

    add_footer_line(slide,
        "\"Кто я\" не должно смешиваться с \"что я знаю\" и \"что я делаю прямо сейчас\"")


# ─── Slide 3: Node — how a knowledge file works ──────────────

def slide_node(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("Node", PINK), (" — файл знаний", TEAL)
    ])
    add_subtitle(slide, "Один файл = одна законченная мысль")

    # File name example (in a light box)
    add_card(slide, 0.5, 1.6, 9.0, 0.65, LIGHT_GRAY, None)
    add_text(slide,
        "Имя файла: structure-enables-navigation-without-reading-everything.md",
        0.7, 1.68, 8.6, 0.5, size=14, color=GRAY)

    # Two columns: YAML | Body
    # Left: YAML
    add_text(slide, "YAML-шапка", 0.7, 2.55, 3.8, 0.35,
             font=FONT_H, size=16, color=PINK, bold=True)

    yaml_lines = [
        ("description:", TEAL, True),
        ("  суть файла в ~150 символов", DARK_BLUE, False),
        ("kind: insight | pattern | decision", TEAL, True),
        ("status: preliminary | active", TEAL, True),
        ("topics: [\"topic-1\", \"topic-2\"]", TEAL, True),
    ]
    add_multiline(slide, yaml_lines, 0.7, 2.95, 4.0, 1.8, size=13, spacing=4)

    # Vertical divider
    add_rect(slide, 4.85, 2.55, 0.02, 2.2, CARD_BORDER)

    # Right: Body
    add_text(slide, "Тело файла", 5.1, 2.55, 4.2, 0.35,
             font=FONT_H, size=16, color=BRIGHT_TEAL, bold=True)

    body_lines = [
        "Самодостаточная проза.",
        "Файл понятен без чтения других файлов.",
        "Wikilinks [[внутри предложений]],",
        "не списком в конце.",
    ]
    add_multiline(slide, body_lines, 5.1, 2.95, 4.2, 1.8, size=14, spacing=6)


# ─── Slide 4: YAML as queryable database ─────────────────────

def slide_yaml_db(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("YAML", PINK), (" как queryable database", TEAL)
    ])
    add_subtitle(slide, "249 markdown-файлов → SQL-запросы без базы данных")

    # Command → Result pairs
    commands = [
        ("grep 'kind: tension'", "→ все противоречия в системе"),
        ("grep 'kind: decision'", "→ все принятые решения"),
        ("grep 'status: preliminary'", "→ всё что ещё не проверено"),
        ("grep 'topics: .*найм'", "→ всё что связано с наймом"),
    ]

    for i, (cmd, result) in enumerate(commands):
        y = 1.65 + i * 0.75

        # Dark code block
        code_card = add_card(slide, 0.5, y, 4.0, 0.55, CODE_BG, None)
        add_text(slide, cmd, 0.7, y + 0.08, 3.6, 0.4,
                 size=15, color=WHITE)

        # Result
        add_text(slide, result, 4.8, y + 0.08, 4.7, 0.4,
                 size=17, color=TEAL)


# ─── Slide 5: Section — How Files Connect ────────────────────

def slide_section_connect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_TEAL_BG
    brand_markers(slide)

    add_text(slide, "Как файлы связаны", 1.0, 1.8, 8.0, 1.5,
             font=FONT_H, size=44, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)

    # Decorative bar
    add_rect(slide, 4.0, 3.5, 2.0, 0.06, TEAL)


# ─── Slide 6: MOC Hierarchy ──────────────────────────────────

def slide_moc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title(slide, "Вертикальная структура: MOC", size=36)
    add_subtitle(slide, "Map of Content — файлы-навигаторы, не содержат знания сами")

    # Stepped hierarchy: each level is a wider bar
    levels = [
        ("Hub MOC", "Точка входа в базу знаний", CODE_BG, WHITE, 7.5),
        ("Domain MOC", "7 доменов", CODE_BG, WHITE, 6.2),
        ("Topic MOC", "Кластеры по 5-20 нодов", SUBTLE_TEAL, DARK_BLUE, 4.8),
        ("Node", "Атомарный файл знания", SUBTLE_TEAL, DARK_BLUE, 3.5),
    ]

    for i, (name, desc, bg, text_color, bar_w) in enumerate(levels):
        y = 1.6 + i * 0.85
        x = 0.5

        # Bar
        bar = add_card(slide, x, y, bar_w, 0.55, bg, None)

        # Name (bold) + description in one line
        tf = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.08), Inches(bar_w - 0.4), Inches(0.4)
        ).text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = name + "  "
        r1.font.name = FONT_H
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = text_color
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = FONT_B
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY if bg == SUBTLE_TEAL else RGBColor(0xA0, 0xA0, 0xA0)

        # Arrow between levels
        if i < len(levels) - 1:
            add_text(slide, "↓", 0.7, y + 0.5, 0.5, 0.35,
                     size=18, color=GRAY, align=PP_ALIGN.CENTER)


# ─── Slide 7: Wikilinks ──────────────────────────────────────

def slide_wikilinks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title(slide, "Горизонтальные связи: wikilinks", size=34)

    # Left: BAD (✗)
    add_card(slide, 0.5, 1.3, 4.1, 2.8, WHITE, CARD_BORDER)
    add_text(slide, "X   Список внизу", 0.7, 1.4, 3.7, 0.35,
             font=FONT_H, size=16, color=PINK, bold=True)

    bad_lines = [
        "## Related:",
        "[[файл-1]], [[файл-2]]",
        "",
        "Агент видит ЧТО связано,",
        "но не понимает ЗАЧЕМ",
    ]
    add_multiline(slide, bad_lines, 0.7, 1.85, 3.7, 2.0, size=14, spacing=4)

    # Right: GOOD (✓)
    add_card(slide, 5.2, 1.3, 4.3, 2.8, WHITE, CARD_BORDER)

    # Green accent bar
    add_rect(slide, 5.2, 1.3, 4.3, 0.06, TEAL)

    add_text(slide, "OK  В прозе", 5.4, 1.4, 3.9, 0.35,
             font=FONT_H, size=16, color=TEAL, bold=True)

    good_lines = [
        "Решение принято после",
        "[[провал-аутсорса-Q2]] — стало",
        "ясно что нужен свой человек.",
        "",
        "Агент видит КОНТЕКСТ связи",
    ]
    add_multiline(slide, good_lines, 5.4, 1.85, 3.9, 2.0, size=14, spacing=4)

    # Articulation test box
    add_card(slide, 0.5, 4.3, 9.0, 0.65, SUBTLE_TEAL, None)
    add_text(slide,
        "Articulation test: \"A связан с B потому что [конкретная причина]\". Голое \"see also\" не допускается.",
        0.7, 4.38, 8.6, 0.5, size=15, color=TEAL)

    add_footer_line(slide,
        "MOC — вертикаль.  Wikilinks — горизонталь.  Вместе — единое связанное пространство.",
        y=5.05)


# ─── Slide 8: Graph, not tree ────────────────────────────────

def slide_graph(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("Граф", PINK), (", не дерево", TEAL)
    ], size=38)

    # Central node: "Ищем CTO с опытом в AI"
    add_card(slide, 0.5, 1.4, 9.0, 0.65, CODE_BG, None)

    tf = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.48), Inches(8.6), Inches(0.5)
    ).text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Ищем CTO с опытом в AI"
    r.font.name = FONT_H
    r.font.size = Pt(22)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Three arrows down
    add_text(slide, "↓", 1.5, 2.1, 0.5, 0.35, size=20, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "↓", 4.75, 2.1, 0.5, 0.35, size=20, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "↓", 8.0, 2.1, 0.5, 0.35, size=20, color=GRAY, align=PP_ALIGN.CENTER)

    # Three domain cards
    domains = [
        ("Найм", "Домен: Команда", "потому что это про людей", TEAL),
        ("Продукт", "Домен: Продукт", "от CTO зависит техническое направление", BRIGHT_TEAL),
        ("Бюджет", "Домен: Финансы", "расход $200K+ в год", DARK_BLUE),
    ]

    for i, (title, domain, reason, accent) in enumerate(domains):
        x = 0.5 + i * 3.15
        y = 2.55

        add_card(slide, x, y, 2.8, 1.6, WHITE, CARD_BORDER)
        add_rect(slide, x, y, 2.8, 0.06, accent)

        add_text(slide, title, x + 0.15, y + 0.2, 2.5, 0.35,
                 font=FONT_H, size=20, color=accent, bold=True)
        add_text(slide, domain, x + 0.15, y + 0.6, 2.5, 0.3,
                 size=12, color=GRAY)
        add_text(slide, reason, x + 0.15, y + 0.95, 2.5, 0.55,
                 size=13, color=DARK_BLUE)

    add_footer_line(slide,
        "В папках — один путь. В графе — три точки доступа к одному факту.",
        y=4.5)


# ─── Slide 9: 5 Layers ───────────────────────────────────────

def slide_five_layers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("5 слоёв", PINK), (" чтения графа", TEAL)
    ], size=38)
    add_subtitle(slide, "От дешёвого к дорогому. Останавливается, когда получен ответ.")

    # Progressive bars (wider = cheaper/faster, narrower = expensive)
    layers = [
        ("1", "Дерево файлов + self/", "мгновенно", 7.5, CODE_BG, WHITE),
        ("2", "YAML descriptions", "секунды", 6.8, CODE_BG, WHITE),
        ("3", "MOC-навигаторы", "быстро", 6.0, CODE_BG, WHITE),
        ("4", "Wikilinks в прозе", "средне", 5.2, SUBTLE_TEAL, DARK_BLUE),
        ("5", "Полное чтение", "дорого", 4.4, SUBTLE_TEAL, DARK_BLUE),
    ]

    for i, (num, name, speed, bar_w, bg, text_color) in enumerate(layers):
        y = 1.55 + i * 0.7
        x = 0.5

        add_card(slide, x, y, bar_w, 0.5, bg, None)

        # Number
        tf = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.05), Inches(0.3), Inches(0.4)
        ).text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.name = FONT_H
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = text_color

        # Name
        add_text(slide, name, x + 0.55, y + 0.05, bar_w - 1.0, 0.4,
                 font=FONT_H, size=17, color=text_color)

        # Speed label (right-aligned outside bar)
        add_text(slide, speed, bar_w + 0.7, y + 0.07, 1.5, 0.35,
                 size=13, color=GRAY)


# ─── Slide 10: Stats — 249 → 5-10 ────────────────────────────

def slide_stats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    # Two big numbers side by side
    # Left: 249
    add_text(slide, "249", 0.5, 0.8, 4.2, 1.5,
             font=FONT_H, size=72, color=PINK, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "файлов в базе", 0.5, 2.2, 4.2, 0.6,
             font=FONT_B, size=28, color=DARK_BLUE,
             align=PP_ALIGN.CENTER)

    # Right: 5-10
    add_text(slide, "5-10", 5.3, 0.8, 4.2, 1.5,
             font=FONT_H, size=72, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "читает агент", 5.3, 2.2, 4.2, 0.6,
             font=FONT_B, size=28, color=DARK_BLUE,
             align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 4.85, 1.0, 0.03, 1.8, CARD_BORDER)

    # Explanation box
    add_card(slide, 0.5, 3.2, 9.0, 1.2, SUBTLE_TEAL, None)
    add_multiline(slide, [
        "Большинство вопросов решается на слоях 2-3.",
        "Плюс параллельный канал — семантический поиск через embeddings,",
        "который находит то, что ссылки пропускают."
    ], 0.7, 3.35, 8.6, 1.0, size=17, spacing=6)


# ─── Slide 11: Skills ────────────────────────────────────────

def slide_skills(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("Скиллы — навигаторы,", TEAL)
    ], y=0.15, h=0.5, size=30)
    add_title_two_colors(slide, [
        ("не контейнеры", PINK)
    ], y=0.55, h=0.5, size=30)

    # Key insight
    add_subtitle(slide, "Скиллы не содержат знания. Содержат инструкции КАК доставать знания из графа.",
                 y=1.05, size=16)

    # Three skill cards
    skills = [
        ("/architect", TEAL, [
            "Сканирует здоровье графа,",
            "находит проблемные места,",
            "выдаёт рекомендации"
        ]),
        ("/reduce", PINK, [
            "Извлекает structured insights,",
            "создаёт новый node,",
            "проверяет YAML"
        ]),
        ("/reflect", BRIGHT_TEAL, [
            "Проходит по wikilinks,",
            "ищет противоречия,",
            "добавляет новые связи"
        ]),
    ]

    for i, (name, accent, lines) in enumerate(skills):
        x = 0.5 + i * 3.15
        y = 1.55

        add_card(slide, x, y, 2.8, 2.2, WHITE, CARD_BORDER)
        add_rect(slide, x, y, 2.8, 0.06, accent)

        add_text(slide, name, x + 0.15, y + 0.2, 2.5, 0.4,
                 font=FONT_H, size=22, color=accent, bold=True)

        add_multiline(slide, lines, x + 0.15, y + 0.7, 2.5, 1.4, size=13, spacing=4)

    # Two types section
    add_rect(slide, 0.5, 3.95, 9.0, 0.02, CARD_BORDER)

    # Plugin-level
    add_text(slide, "Plugin-level (10)", 0.5, 4.1, 4.0, 0.3,
             font=FONT_H, size=14, color=TEAL, bold=True)
    add_text(slide, "НАД графом: setup, architect, health", 0.5, 4.38, 4.0, 0.25,
             size=11, color=GRAY)

    # Generated
    add_text(slide, "Generated (16)", 5.5, 4.1, 4.0, 0.3,
             font=FONT_H, size=14, color=PINK, bold=True)
    add_text(slide, "ВНУТРИ графа: reduce, reflect, reweave", 5.5, 4.38, 4.0, 0.25,
             size=11, color=GRAY)

    add_footer_line(slide,
        "Граф — мозг (знания и связи).  Скиллы — навыки (что мозг умеет делать).",
        y=4.75)


# ─── Slide 12: Hooks ─────────────────────────────────────────

def slide_hooks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    brand_markers(slide)

    add_title_two_colors(slide, [
        ("Хуки", PINK), (" — рефлексы системы", TEAL)
    ], size=36)
    add_subtitle(slide, "Enforcement, не intelligence. Проверяют и сигнализируют, не принимают решения.",
                 size=16)

    hooks = [
        ("session-orient", "при старте сессии", TEAL, "Утренний брифинг",
         "Загружает контекст + сигналы обслуживания:\n«inbox переполнен», «5+ противоречий», «methodology устарел»"),
        ("write-validate", "при записи файла", PINK, "Фейс-контроль",
         "Проверяет YAML-шапку: description, topics, разделители.\nNon-blocking — предупреждает, не блокирует."),
        ("auto-commit", "при записи файла", BRIGHT_TEAL, "Страховка",
         "Автоматический git commit.\nАсинхронный, незаметный. Всегда можно откатить."),
    ]

    for i, (name, trigger, accent, label, desc) in enumerate(hooks):
        y = 1.65 + i * 1.1

        # Left: name + trigger
        add_text(slide, name, 0.6, y, 3.0, 0.35,
                 font=FONT_H, size=19, color=accent, bold=True)

        tf = slide.shapes.add_textbox(
            Inches(0.6), Inches(y + 0.35), Inches(3.0), Inches(0.55)
        ).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = trigger
        r1.font.name = FONT_B
        r1.font.size = Pt(11)
        r1.font.color.rgb = GRAY
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = FONT_B
        r2.font.size = Pt(14)
        r2.font.color.rgb = DARK_BLUE

        # Right: description
        add_multiline(slide, desc.split("\n"),
                      4.2, y + 0.05, 5.3, 0.9, size=13, spacing=4)

        # Separator
        if i < len(hooks) - 1:
            add_rect(slide, 0.6, y + 1.0, 8.8, 0.01, CARD_BORDER)


# ─── Slide 13: Closing Quote ─────────────────────────────────

def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_TEAL_BG
    brand_markers(slide)

    # Large quotation mark
    add_text(slide, "\u201C", 0.8, 0.5, 1.0, 1.5,
             font=FONT_H, size=80, color=BRIGHT_TEAL, bold=True)

    # Quote
    add_text(slide, "Граф — мозг.\nСкиллы — навыки.\nХуки — рефлексы.",
             1.5, 1.5, 7.0, 2.5,
             font=FONT_B, size=32, color=DARK_BLUE, italic=True)

    # Attribution
    add_text(slide, "— arscontexta", 1.5, 4.0, 7.0, 0.5,
             size=16, color=GRAY, align=PP_ALIGN.RIGHT)


# ─── Main ─────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slide_title(prs)
    slide_three_spaces(prs)
    slide_node(prs)
    slide_yaml_db(prs)
    slide_section_connect(prs)
    slide_moc(prs)
    slide_wikilinks(prs)
    slide_graph(prs)
    slide_five_layers(prs)
    slide_stats(prs)
    slide_skills(prs)
    slide_hooks(prs)
    slide_closing(prs)

    output = "/Users/viktorsolomonik/Projects/svaib/_inbox/arscontexta-v3.pptx"
    prs.save(output)
    print(f"Created: {output} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
