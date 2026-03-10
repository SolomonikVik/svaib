#!/usr/bin/env python3
"""Analyze two PPTX files and compare their structure."""
import sys
sys.path.insert(0, "/Users/viktorsolomonik/Projects/svaib/.venv/lib/python3.14/site-packages")

from pptx import Presentation

def safe_color(font):
    try:
        if font.color and font.color.rgb:
            return str(font.color.rgb)
    except:
        pass
    try:
        if font.color and font.color.theme_color:
            return f"theme:{font.color.theme_color}"
    except:
        pass
    return "inherit"

def analyze_pptx(path, label):
    prs = Presentation(path)
    print(f"\n{'='*60}")
    print(f"  {label}")
    print(f"  Slides: {len(prs.slides)}, Size: {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f} inches")
    print(f"{'='*60}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            pos = f"x={shape.left/914400:.2f} y={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}"

            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        font_info = ""
                        if para.runs:
                            r = para.runs[0]
                            sz = int(r.font.size/12700) if r.font.size else "?"
                            font_info = f" [{r.font.name or '?'}/{sz}pt/{safe_color(r.font)}]"
                        texts.append(f"'{para.text[:60]}'{font_info}")
                if texts:
                    print(f"  [{j}] TEXT {pos}")
                    for t in texts:
                        print(f"      {t}")
            else:
                try:
                    if hasattr(shape, 'image'):
                        img = shape.image
                        print(f"  [{j}] IMAGE {pos} ({img.content_type}, {len(img.blob)//1024}KB)")
                        continue
                except:
                    pass
                print(f"  [{j}] SHAPE {pos}")

analyze_pptx("/Users/viktorsolomonik/Projects/svaib/_inbox/arscontexta.pptx", "OUT-OF-BOX (score: 4)")
print("\n\n" + "#"*80 + "\n")
analyze_pptx("/Users/viktorsolomonik/Projects/svaib/_inbox/arscontexta-context-engineering.pptx", "OUR SKILL (score: 1-2)")
